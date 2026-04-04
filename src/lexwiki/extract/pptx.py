"""PPTX to markdown extraction using python-pptx."""

from __future__ import annotations

import io
from pathlib import Path


def extract_pptx(source: Path) -> tuple[str, int]:
    """Convert PPTX to structured markdown.

    Each slide becomes a ## heading with text content, tables, and speaker notes.
    Returns (markdown_text, slide_count).
    """
    from pptx import Presentation

    blob = source.read_bytes()
    prs = Presentation(io.BytesIO(blob))
    lines = []

    for i, slide in enumerate(prs.slides, 1):
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text
        lines.append(f"## Slide {i}" + (f": {title}" if title else ""))
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text and text != title:
                        lines.append(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    lines.append("| " + " | ".join(cells) + " |")
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append(f"\n> Notes: {notes}")
        lines.append("")

    slide_count = len(prs.slides)
    md = "\n".join(lines) if lines else "[Empty presentation]"
    return md, slide_count
